from flask import Flask, request, jsonify, send_from_directory, render_template
import os
import uuid
import json
import time
from audio_utils import standardize_audio, check_audio_duration  # 自定义音频处理工具
import pyttsx3  # 文本转语音
from pptx import Presentation  # 用于解析 PPT 文件内容

app = Flask(__name__)

# ---------- 文件路径配置 ----------
UPLOAD_FOLDER = 'uploads'                 # 所有上传内容统一放在 uploads 下
PPT_FOLDER = os.path.join(UPLOAD_FOLDER, 'pptx')  # PPT 文件子目录
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(PPT_FOLDER, exist_ok=True)

# ---------- 元数据存储路径 ----------
METADATA_PATH = os.path.join(UPLOAD_FOLDER, 'audio_metadata.json')  # 音频样本信息
TEXT_TASK_PATH = os.path.join(UPLOAD_FOLDER, 'text_tasks.json')     # 文本任务信息

# ---------- 工具函数 ----------

# 保存新上传的音频元信息（追加存储）
def save_metadata(metadata):
    if os.path.exists(METADATA_PATH):
        with open(METADATA_PATH, 'r') as f:
            data = json.load(f)
    else:
        data = []
    data.append(metadata)
    with open(METADATA_PATH, 'w') as f:
        json.dump(data, f, indent=2)

# 加载所有声音样本信息
def load_all_voices():
    if os.path.exists(METADATA_PATH):
        with open(METADATA_PATH, 'r') as f:
            return json.load(f)
    return []

# ---------- 欢迎页及主界面 ----------

@app.route("/", methods=["GET"])
def welcome():
    """
    欢迎界面
    """
    return render_template("welcome.html")

@app.route("/main", methods=["GET"])
def index():
    """
    主业务界面
    """
    return render_template("index.html")

# ---------- 接口定义区 ----------

# ① 上传教师音频样本
@app.route("/upload_audio", methods=["POST"])
def upload_audio():
    """
    上传音频文件，并自动标准化（转为16kHz单声道wav），记录音频元数据。
    """
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files['file']
    filename = file.filename
    original_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(original_path)

    try:
        # 标准化音频
        processed_path = standardize_audio(original_path)
        duration = check_audio_duration(processed_path)

        # 为该音频样本生成唯一ID
        voice_id = str(uuid.uuid4())
        metadata = {
            "id": voice_id,
            "filename": os.path.basename(processed_path),
            "duration": duration,
            "upload_time": time.strftime('%Y-%m-%d %H:%M:%S'),
            "path": processed_path
        }
        save_metadata(metadata)
        return jsonify({"message": "上传成功", "id": voice_id}), 200

    except Exception as e:
        return jsonify({"error": str(e)}), 500

# ② 获取所有上传的声音样本列表
@app.route("/list_voices", methods=["GET"])
def list_voices():
    """列出所有声音样本信息"""
    return jsonify(load_all_voices()), 200

# ③ 提交教学文本任务（将文本与声音样本绑定）
@app.route("/submit_text_task", methods=["POST"])
def submit_text_task():
    """
    提交一个新的文本转语音任务
    """
    data = request.get_json()
    text = data.get("text", "").strip()
    voice_id = data.get("voice_id")

    # 参数校验
    if not text or not voice_id:
        return jsonify({"error": "text 和 voice_id 不能为空"}), 400

    if len(text) < 800 or len(text) > 2000:
        return jsonify({"error": "文本长度应在800到2000字符之间"}), 400

    voices = load_all_voices()
    voice_exists = any(v["id"] == voice_id for v in voices)
    if not voice_exists:
        return jsonify({"error": "无效的 voice_id"}), 400

    # 创建任务对象
    task_id = str(uuid.uuid4())
    task_record = {
        "task_id": task_id,
        "voice_id": voice_id,
        "text": text,
        "submit_time": time.strftime('%Y-%m-%d %H:%M:%S'),
        "status": "pending"
    }

    # 保存任务到任务列表
    if os.path.exists(TEXT_TASK_PATH):
        with open(TEXT_TASK_PATH, 'r') as f:
            tasks = json.load(f)
    else:
        tasks = []

    tasks.append(task_record)
    with open(TEXT_TASK_PATH, 'w') as f:
        json.dump(tasks, f, indent=2)

    return jsonify({"message": "文本任务已提交", "task_id": task_id}), 200

# ④ 根据任务生成语音（使用 pyttsx3 占位合成）
@app.route("/generate_audio", methods=["POST"])
def generate_audio():
    """
    根据任务ID，使用 pyttsx3 合成语音，生成 wav 文件
    """
    data = request.get_json()
    task_id = data.get("task_id")

    if not task_id:
        return jsonify({"error": "缺少 task_id"}), 400

    if not os.path.exists(TEXT_TASK_PATH):
        return jsonify({"error": "任务列表为空"}), 400

    with open(TEXT_TASK_PATH, 'r') as f:
        tasks = json.load(f)

    task = next((t for t in tasks if t["task_id"] == task_id), None)
    if not task:
        return jsonify({"error": "未找到任务"}), 404

    text = task["text"]
    output_filename = f"{task_id}_output.wav"
    output_path = os.path.join(UPLOAD_FOLDER, output_filename)

    try:
        # 使用 pyttsx3 将文本转换为语音
        engine = pyttsx3.init()
        engine.setProperty('rate', 150)
        engine.save_to_file(text, output_path)
        engine.runAndWait()

        task["status"] = "completed"
        task["output_audio"] = output_filename

        with open(TEXT_TASK_PATH, 'w') as f:
            json.dump(tasks, f, indent=2)

        return jsonify({
            "message": "音频生成成功",
            "audio_file": output_filename,
            "download_url": f"/get_audio/{task_id}"
        }), 200

    except Exception as e:
        return jsonify({"error": str(e)}), 500

# ⑤ 下载任务生成的语音
@app.route("/get_audio/<task_id>", methods=["GET"])
def get_audio(task_id):
    """
    提供音频下载（根据任务ID）
    """
    if not os.path.exists(TEXT_TASK_PATH):
        return jsonify({"error": "任务记录不存在"}), 404

    with open(TEXT_TASK_PATH, 'r') as f:
        tasks = json.load(f)

    task = next((t for t in tasks if t["task_id"] == task_id), None)
    if not task or task.get("status") != "completed":
        return jsonify({"error": "任务未完成或未找到"}), 404

    filename = task.get("output_audio")
    if not filename or not os.path.exists(os.path.join(UPLOAD_FOLDER, filename)):
        return jsonify({"error": "音频文件不存在"}), 404

    return send_from_directory(UPLOAD_FOLDER, filename, as_attachment=True)

# ⑥ 获取所有任务列表
@app.route("/list_text_tasks", methods=["GET"])
def list_text_tasks():
    """
    获取所有文本转语音任务
    """
    if not os.path.exists(TEXT_TASK_PATH):
        return jsonify([]), 200

    with open(TEXT_TASK_PATH, 'r') as f:
        tasks = json.load(f)
    return jsonify(tasks), 200

# ⑦ 删除任务（含音频）
@app.route("/delete_task/<task_id>", methods=["DELETE"])
def delete_task(task_id):
    """
    删除某个任务，同时删除对应生成的音频文件
    """
    if not os.path.exists(TEXT_TASK_PATH):
        return jsonify({"error": "任务列表不存在"}), 404

    with open(TEXT_TASK_PATH, 'r') as f:
        tasks = json.load(f)

    new_tasks = [t for t in tasks if t["task_id"] != task_id]

    # 同时删除已生成的音频文件
    for t in tasks:
        if t["task_id"] == task_id and t.get("output_audio"):
            path = os.path.join(UPLOAD_FOLDER, t["output_audio"])
            if os.path.exists(path):
                os.remove(path)

    with open(TEXT_TASK_PATH, 'w') as f:
        json.dump(new_tasks, f, indent=2)

    return jsonify({"message": f"任务 {task_id} 已删除"}), 200

# ⑧ 上传 PPT 文件并提取文本
@app.route("/upload_ppt", methods=["POST"])
def upload_ppt():
    """
    上传 PPT 文件，并解析每页文字内容
    """
    if 'file' not in request.files:
        return jsonify({"error": "没有上传文件"}), 400

    file = request.files['file']
    filename = file.filename

    if not filename.lower().endswith('.pptx'):
        return jsonify({"error": "只支持 .pptx 文件"}), 400

    ppt_id = str(uuid.uuid4())
    save_path = os.path.join(PPT_FOLDER, f"{ppt_id}.pptx")
    file.save(save_path)

    try:
        # 解析 ppt 文件，提取每页文字内容
        prs = Presentation(save_path)
        slide_texts = []
        for slide in prs.slides:
            content = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content.append(shape.text)
            slide_texts.append('\n'.join(content).strip())

        return jsonify({
            "ppt_id": ppt_id,
            "slides": slide_texts,
            "slide_count": len(slide_texts)
        }), 200

    except Exception as e:
        return jsonify({"error": f"PPT解析失败: {str(e)}"}), 500

# ---------- 启动服务器 ----------
if __name__ == "__main__":
    app.run(debug=True)
